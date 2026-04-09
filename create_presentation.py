"""MindTune 10분 발표 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상 팔레트
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF7)
DEEP_PURPLE = RGBColor(0x2D, 0x1B, 0x69)
MINDTUNE_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
PLACEHOLDER_BG = RGBColor(0xE5, 0xE7, 0xEB)

FONT_KR = "맑은 고딕"
FONT_CODE = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_accent_bar(slide, top=Inches(0), color=MINDTUNE_PURPLE):
    """상단 악센트 바"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_text(slide, text, left=Inches(0.8), top=Inches(0.5), width=Inches(11.5), size=28, color=DEEP_PURPLE, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_KR
    return txBox


def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.5), width=Inches(5), size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_KR
        p.space_after = Pt(8)
    return txBox


def add_placeholder(slide, text, left, top, width, height):
    """캡처 사진 플레이스홀더"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_BG
    shape.line.color.rgb = MEDIUM_GRAY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_KR


def add_table(slide, rows, cols, data, left, top, width, height):
    """테이블 추가"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = FONT_KR
                if i == 0:  # 헤더
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MINDTUNE_PURPLE
    return table


def add_as_is_to_be_box(slide, as_is_text, to_be_text, left, top, width):
    """AS-IS / TO-BE 비교 박스"""
    half = width / 2 - Inches(0.1)
    # AS-IS
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, half, Inches(1.2))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RED_BG
    shape1.line.fill.background()
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "AS-IS"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = RED
    p1.font.name = FONT_KR
    for line in as_is_text.split("\n"):
        p = tf1.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_KR

    # TO-BE
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + half + Inches(0.2), top, half, Inches(1.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN_BG
    shape2.line.fill.background()
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "TO-BE"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = GREEN
    p2.font.name = FONT_KR
    for line in to_be_text.split("\n"):
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_KR


# ========== 슬라이드 1: 표지 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_accent_bar(slide)
add_title_text(slide, "MindTune", top=Inches(2.2), size=48, color=MINDTUNE_PURPLE)
add_body_text(slide, "AI 기반 음악 심리 케어 에이전트", top=Inches(3.2), size=24, color=MEDIUM_GRAY)
add_body_text(slide, "Tool-Calling ReAct Agent  |  하이브리드 검색 (벡터 + BM25)  |  동질 원리 플레이리스트", top=Inches(4.0), size=16, color=MEDIUM_GRAY)
add_body_text(slide, "2026.04", top=Inches(5.5), size=16, color=MEDIUM_GRAY)
# 하단 악센트 바
add_accent_bar(slide, top=Inches(7.42))

# ========== 슬라이드 2: 왜 음악 심리 케어인가? ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "왜 음악 심리 케어인가?")
add_body_text(slide,
    "회사 도메인 연결\n"
    "  - 의료AI + 에이전틱 AI + 멀티모달 데이터 플랫폼\n"
    "  - ACTIC LLM 운영 솔루션 → 에이전트 기반 AX\n\n"
    "데이터셋의 의도\n"
    "  - 'Songs, Ragas, Mental Health Classification'\n"
    "  - Mental_Health_Label 컬럼에 Anxiety, Depression, PTSD 등\n\n"
    "음악치료 임상 근거\n"
    "  - 코크란 체계적 고찰: 우울 증상 감소 효과 입증\n"
    "  - 음악 청취 시 도파민 분비 → 정서적 안정",
    width=Inches(11))

# ========== 슬라이드 3: 음악치료 + 동질 원리 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "음악치료와 동질 원리 (ISO Principle)")
add_body_text(slide,
    "동질 원리란?\n"
    "  - 현재 감정과 유사한 음악에서 시작\n"
    "  - 점진적으로 목표 감정으로 이동\n"
    "  - 감정 변화를 유도하는 플레이리스트 구성\n\n"
    "MindTune 적용\n"
    "  - 2D 감정 공간 (Valence + Energy)\n"
    "  - 유클리드 거리 기반 곡 순서 배열\n"
    "  - 예: 우울(V=0.2, E=0.3) → 중립(V=0.5, E=0.5) → 긍정(V=0.7, E=0.6)",
    width=Inches(6))
add_placeholder(slide, "[동질 원리 다이어그램]\n현재 감정 → 중립 → 목표 감정\nvalence/energy 2D 좌표계",
    Inches(7.5), Inches(1.5), Inches(5), Inches(4.5))

# ========== 슬라이드 4: EDA 데이터 개요 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "데이터셋 EDA: 32,833곡 분석")
add_body_text(slide,
    "데이터 정제\n"
    "  - 원본 32,833곡 → track_id 중복 제거 → track_name+artist 중복 제거\n"
    "  - 최종 26,229곡 (결측값 1건 추가 제거)\n\n"
    "핵심 발견\n"
    "  - 6개 장르 균등 분포 (EDM/Rap/Pop/R&B/Latin/Rock)\n"
    "  - MH 라벨 극심한 불균형: Depression 138곡, PTSD 38곡\n"
    "  - 11개 오디오 피처 + Mental Health Label 활용",
    width=Inches(5.5))
add_placeholder(slide, "[EDA 캡처 삽입]\n레이더차트 / MH 라벨 분포 바차트 / 감정 공간 산점도",
    Inches(7), Inches(1.5), Inches(5.5), Inches(5))

# ========== 슬라이드 5: EDA → 설계 의사결정 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "EDA 인사이트 → 설계 의사결정")
add_body_text(slide,
    "MH 라벨 불균형 (Depression 138곡)\n"
    "  → 라벨 필터만으론 부족 → 오디오 피처 기반 검색으로 보완\n\n"
    "감정-피처 상관관계 (Mann-Whitney U 검정)\n"
    "  → 감정별 valence/energy/acousticness 유의차 확인\n"
    "  → 감정-피처 매핑 가이드 설계 (7개 감정 × 4개 피처)\n\n"
    "장르별 피처 프로파일\n"
    "  → 장르 간 코사인 유사도 분석 → 상황 프리셋 7종 설계\n\n"
    "유사곡 시뮬레이션\n"
    "  → 벡터 단독 검색 한계 발견 → 하이브리드 검색 도입 결정",
    width=Inches(11))

# ========== 슬라이드 6: 시스템 아키텍처 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "시스템 아키텍처: 3-Agent 파이프라인")
add_body_text(slide,
    "Gate Agent (라우터 + 안전장치)\n"
    "  - 키워드 규칙 → LLM 분류 → 체이닝 폴백\n"
    "  - 위험 신호 즉시 감지 (LLM 호출 생략)\n\n"
    "Music Agent (ReAct 루프)\n"
    "  - 6개 도구 자율 선택 (tool_choice=auto)\n"
    "  - 최대 7회 반복, Thinking Mode 동적 전환\n"
    "  - Self-Verification + 폴백 응답 생성\n\n"
    "Care Agent (공감 응답)\n"
    "  - intent별 케어 강도 차별화\n"
    "  - 출력 가드레일 (의료 표현 필터링)\n\n"
    "  + Spotify 검색 병렬 실행",
    width=Inches(5.5))
add_placeholder(slide, "[시스템 아키텍처 다이어그램]\nGate → Music Agent(ReAct) → Care + Spotify 병렬\n6개 도구 목록 포함",
    Inches(7), Inches(1.5), Inches(5.5), Inches(5))

# ========== 슬라이드 7: 하이브리드 검색 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "하이브리드 검색 엔진: 벡터 + BM25 + 피처")
add_body_text(slide,
    "3중 검색 구조\n"
    "  1. 벡터 검색 (sentence-transformers → ChromaDB)\n"
    "     의미적 유사도 기반\n"
    "  2. BM25 키워드 검색 (rank_bm25)\n"
    "     정확한 키워드 매칭 보완\n"
    "  3. 피처 검색 (Pandas + 코사인 유사도)\n"
    "     수치 범위 필터링 + 7개 피처 벡터 랭킹\n\n"
    "하이브리드 점수\n"
    "  final = 0.6 × vector + 0.4 × BM25\n\n"
    "리랭커 가중치 (intent별 차별화)\n"
    "  emotion:  피처 60% | 벡터 30% | 인기도 10%\n"
    "  situation: 피처 50% | 벡터 40% | 인기도 10%\n"
    "  similar:  벡터 50% | 피처 40% | 인기도 10%",
    width=Inches(5.5))
add_placeholder(slide, "[하이브리드 검색 다이어그램]\n사용자 쿼리 → 벡터/BM25/피처 → 리랭커 → 최종 결과",
    Inches(7), Inches(1.5), Inches(5.5), Inches(5))

# ========== 슬라이드 8: 데모 - 감정 기반 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "데모 1: 감정 기반 추천 — \"오늘 너무 우울해\"")
add_body_text(slide,
    "처리 흐름\n"
    "  1. Gate: emotion 분류\n"
    "  2. Music: search_by_features(valence=[0.1,0.3])\n"
    "     → rerank → build_iso_playlist\n"
    "  3. Care: 공감 메시지 + 동질 원리 설명\n"
    "  4. Spotify: 임베드 플레이어 제공\n\n"
    "동질 원리 적용\n"
    "  현재(V=0.15, E=0.25) → 중립 → 긍정",
    width=Inches(5))
add_placeholder(slide, "[캡처 삽입]\n감정 추천 결과\n곡 카드 + Spotify 플레이어",
    Inches(6.5), Inches(1.2), Inches(6), Inches(3.5))
add_placeholder(slide, "[캡처 삽입] 사이드바 감정 변화 그래프 (plotly scatter)",
    Inches(6.5), Inches(5), Inches(6), Inches(2))

# ========== 슬라이드 9: 데모 - 상황 맞춤 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "데모 2: 상황 맞춤 BGM — \"카페에서 책 읽을 때\"")
add_body_text(slide,
    "처리 흐름\n"
    "  1. Gate: situation 분류 (키워드 규칙, LLM 생략)\n"
    "  2. Music: 프리셋 매칭 → search_by_description\n"
    "     → search_by_features → rerank\n"
    "  3. Care: 가벼운 상황 코멘트\n\n"
    "대기 중 UX\n"
    "  데이터셋 서브장르 랜덤 트렌딩 Top 3\n"
    "  Spotify 임베드로 바로 재생 가능",
    width=Inches(5))
add_placeholder(slide, "[캡처 삽입]\n상황 BGM 추천 결과 + 트렌딩 카드",
    Inches(6.5), Inches(1.2), Inches(6), Inches(5.5))

# ========== 슬라이드 10: 데모 - 유사곡 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "데모 3: 유사곡 탐색 — \"Coldplay Yellow 같은 곡\"")
add_body_text(slide,
    "처리 흐름\n"
    "  1. Gate: similar 분류 (키워드 규칙)\n"
    "  2. Music: lookup_song(Yellow)\n"
    "     → search_by_features + search_by_description\n"
    "     → rerank(intent=similar, vector 50%)\n"
    "  3. Spotify 연동\n\n"
    "에이전트 로그 (토글)\n"
    "  각 답변 아래 도구 호출 순서/파라미터 확인",
    width=Inches(5))
add_placeholder(slide, "[캡처 삽입]\n유사곡 결과 + 에이전트 로그 토글 펼친 모습",
    Inches(6.5), Inches(1.2), Inches(6), Inches(5.5))

# ========== 슬라이드 11: Troubleshooting - 검색 품질 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "Troubleshooting: 검색 품질 개선")
add_as_is_to_be_box(slide,
    "임베딩: energy/valence 2개만\n→ 벡터 검색이 다른 피처 구분 불가",
    "11개 피처 5단계 키워드 + 수치값\n→ 모든 피처가 벡터 공간에 반영",
    Inches(0.8), Inches(1.5), Inches(11.5))
add_as_is_to_be_box(slide,
    "sentence-transformers 단독\n→ 키워드 정확 매칭 부족",
    "BM25 하이브리드 도입\n→ 0.6*벡터 + 0.4*BM25 융합",
    Inches(0.8), Inches(3.0), Inches(11.5))
add_as_is_to_be_box(slide,
    "popularity 가중치 0.1 고정\n→ '유명한 곡' 요청 미반영",
    "키워드 감지 시 동적 상향 (0.1→0.4)\n→ '유명', '인기', '히트' 등 감지",
    Inches(0.8), Inches(4.5), Inches(11.5))
add_as_is_to_be_box(slide,
    "ChromaDB 메타데이터 3개 피처\n→ 필터링 제한",
    "11개 오디오 피처 전체 저장\n→ 모든 피처로 메타데이터 필터 가능",
    Inches(0.8), Inches(6.0), Inches(11.5))

# ========== 슬라이드 12: Troubleshooting - 성능 & 안정성 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "Troubleshooting: 성능 & 안정성 개선")
add_as_is_to_be_box(slide,
    "Gate Agent: 매번 LLM 호출\n→ 분류에 ~1초 소요",
    "키워드 규칙 사전 분류\n→ similar/situation은 LLM 호출 생략",
    Inches(0.8), Inches(1.5), Inches(11.5))
add_as_is_to_be_box(slide,
    "대화 전체를 LLM에 전달\n→ 토큰 낭비 + 컨텍스트 오염",
    "Sliding Window 4턴(8메시지)\n+ 5턴마다 자동 요약 → ~90% 토큰 절약",
    Inches(0.8), Inches(3.0), Inches(11.5))
add_as_is_to_be_box(slide,
    "Care Agent: enable_thinking=False\n→ LLM content=None 반환",
    "옵션 제거 + <think> 코드 처리\n+ reasoning_content 폴백",
    Inches(0.8), Inches(4.5), Inches(11.5))
add_placeholder(slide, "[캡처 삽입] intent별 진행 메시지 변화\n\"감정에 어울리는 곡을 찾고 있어요...\"",
    Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.2))

# ========== 슬라이드 13: 안전장치 + 품질 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "안전장치 & 품질 관리")
add_body_text(slide,
    "입력 안전 체크\n"
    "  - 9개 위험 키워드 사전 감지 (LLM 호출 없이)\n"
    "  - 긴급: 정신건강 위기상담 1577-0199, 자살예방 1393\n"
    "  - 범위 외 질문 → 서비스 안내\n\n"
    "출력 가드레일\n"
    "  - 금지 표현 6종 자동 필터링 (치료, 진단, 처방 등)\n"
    "  - Care Agent 스킵 시에도 orchestrator에서 보장\n\n"
    "할루시네이션 방지\n"
    "  - Self-Verification: 추천 곡 DB 존재 확인\n"
    "  - Grounded Generation: 데이터셋 외 곡 생성 금지\n\n"
    "평가 시스템\n"
    "  - 자동: 할루시네이션/Intent 정확도/Spotify 매칭/Latency\n"
    "  - 수동: Streamlit 평가 UI (1~5점)",
    width=Inches(5.5))
add_placeholder(slide, "[캡처 삽입]\n위기 감지 응답\n(\"죽고 싶어\" → 상담 전화 안내)",
    Inches(7), Inches(1.5), Inches(5.5), Inches(5))

# ========== 슬라이드 14: 한계점 + 확장 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "한계점 & 확장 방향")
add_body_text(slide,
    "정직한 한계점\n"
    "  - 4B 소형 모델의 Tool Calling 정확도 한계\n"
    "  - Mental Health Label 불균형 (Depression 138곡)\n"
    "  - Spotify 검색 매칭률 한계 (데이터셋 곡 ≠ Spotify)\n"
    "  - 전문 음악치료를 대체하지 않는 보조 도구\n\n"
    "확장 방향\n"
    "  - PHR(개인건강기록) 연동: 수면/운동 데이터 기반 맞춤 추천\n"
    "  - ACTIC 플랫폼 배포: 에이전트별 독립 스케일링\n"
    "  - MCP 프로토콜 도입: Apple Music, YouTube Music 추가 연동\n"
    "  - A/B 테스트: 리랭커 가중치, 프롬프트 자동 실험",
    width=Inches(11))

# ========== 슬라이드 15: 마무리 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_title_text(slide, "감사합니다", top=Inches(2.5), size=44, color=MINDTUNE_PURPLE)
add_body_text(slide,
    "기술 스택\n"
    "Python  |  Streamlit  |  OpenAI-compatible API  |  ChromaDB + sentence-transformers + BM25\n"
    "Spotify Web API  |  SQLite  |  Plotly  |  Docker",
    top=Inches(3.8), size=16, color=MEDIUM_GRAY, width=Inches(11))
add_accent_bar(slide, top=Inches(7.42))

# ========== 저장 ==========
output_path = "/home/joon/music-chatbot/presentation.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path} ({len(prs.slides)}장)")
